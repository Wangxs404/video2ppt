#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Video2PPT - 将视频转换为PowerPoint演示文稿
从视频中提取关键帧，生成PPT
"""

import os
import sys
import cv2
import argparse
import logging
from pathlib import Path
from typing import List, Tuple
import numpy as np
from datetime import datetime

try:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError as e:
    print(f"错误: 缺少必要的库 - {e}")
    print("请运行: pip install -r requirements.txt")
    sys.exit(1)


# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


class Video2PPT:
    """视频转PPT转换器"""
    
    def __init__(self, video_path: str, output_path: str = None, fps_interval: int = 1):
        """
        初始化转换器
        
        Args:
            video_path: 视频文件路径
            output_path: 输出PPT文件路径
            fps_interval: 每隔多少秒提取一帧
        """
        self.video_path = video_path
        self.fps_interval = fps_interval
        self.frames: List[str] = []  # 存储帧路径
        self.frames_dir = None
        
        if not os.path.exists(video_path):
            raise FileNotFoundError(f"视频文件不存在: {video_path}")
        
        if output_path is None:
            base_name = Path(video_path).stem
            output_path = f"{base_name}_output.pptx"
        
        self.output_path = output_path
        logger.info(f"初始化转换器: {video_path} -> {output_path}")
    
    def extract_frames(self) -> None:
        """从视频中提取关键帧"""
        logger.info("开始提取视频帧...")
        
        # 创建临时目录存储帧
        self.frames_dir = "temp_frames"
        os.makedirs(self.frames_dir, exist_ok=True)
        
        # 打开视频文件
        cap = cv2.VideoCapture(self.video_path)
        
        if not cap.isOpened():
            raise ValueError(f"无法打开视频文件: {self.video_path}")
        
        fps = cap.get(cv2.CAP_PROP_FPS)
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = total_frames / fps if fps > 0 else 0
        
        logger.info(f"视频信息 - FPS: {fps:.2f}, 总帧数: {total_frames}, 时长: {duration:.2f}秒")
        
        frame_interval = int(fps * self.fps_interval)
        frame_count = 0
        extracted_count = 0
        
        while True:
            ret, frame = cap.read()
            
            if not ret:
                break
            
            # 按间隔提取帧
            if frame_count % frame_interval == 0:
                frame_path = os.path.join(self.frames_dir, f"frame_{extracted_count:04d}.jpg")
                cv2.imwrite(frame_path, frame)
                self.frames.append(frame_path)
                extracted_count += 1
                
                if extracted_count % 10 == 0:
                    logger.info(f"已提取 {extracted_count} 帧")
            
            frame_count += 1
        
        cap.release()
        logger.info(f"完成帧提取，共提取 {extracted_count} 帧")
    
    def generate_ppt(self) -> None:
        """生成PowerPoint演示文稿"""
        logger.info("开始生成PowerPoint演示文稿...")
        
        if not self.frames:
            logger.error("没有可用的帧数据")
            return
        
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # 添加标题幻灯片
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Video2PPT"
        subtitle.text = f"视频转换时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n" \
                       f"源文件: {os.path.basename(self.video_path)}"
        
        # 添加每一帧的幻灯片
        blank_slide_layout = prs.slide_layouts[6]  # 空白布局
        
        for idx, frame_path in enumerate(self.frames, 1):
            logger.info(f"处理第 {idx}/{len(self.frames)} 张幻灯片...")
            
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加帧图像，占据整个幻灯片页面（边缘对齐）
            left = Inches(0)
            top = Inches(0)
            width = Inches(10)     # 幻灯片宽度
            height = Inches(7.5)   # 幻灯片高度
            pic = slide.shapes.add_picture(frame_path, left, top, width=width, height=height)
        
        # 保存演示文稿
        prs.save(self.output_path)
        logger.info(f"PowerPoint已保存: {self.output_path}")
    
    def cleanup(self) -> None:
        """清理临时文件"""
        if self.frames_dir and os.path.exists(self.frames_dir):
            import shutil
            shutil.rmtree(self.frames_dir)
            logger.info("临时文件已清理")
    
    def convert(self) -> None:
        """执行完整的转换流程"""
        try:
            self.extract_frames()
            self.generate_ppt()
            logger.info("转换完成！")
        except Exception as e:
            logger.error(f"转换过程中出错: {e}")
            raise
        finally:
            self.cleanup()


def main():
    """主函数"""
    parser = argparse.ArgumentParser(
        description='将视频文件转换为PowerPoint演示文稿',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  python video2ppt.py input_video.mp4
  python video2ppt.py input_video.mp4 -o output.pptx
  python video2ppt.py input_video.mp4 -i 2  (每2秒提取一帧)
        """
    )
    
    parser.add_argument('video', help='输入视频文件路径')
    parser.add_argument('-o', '--output', help='输出PPT文件路径（默认为input_output.pptx）')
    parser.add_argument('-i', '--interval', type=int, default=1,
                       help='帧提取间隔（秒），默认为1秒')
    
    args = parser.parse_args()
    
    try:
        converter = Video2PPT(args.video, args.output, args.interval)
        converter.convert()
    except Exception as e:
        logger.error(f"错误: {e}")
        sys.exit(1)


if __name__ == '__main__':
    main()
